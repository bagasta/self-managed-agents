"""
File text extractor for RAG document upload.

Supported formats
-----------------
.txt / .md          UTF-8 decode
.pdf                Mistral OCR API (mistral-ocr-latest)
.docx               python-docx (paragraphs + tables)
.pptx               python-pptx (slide text + speaker notes)

After extraction the text is chunked so each chunk fits the embedding
model's context window (~512 tokens ≈ 1200 chars) and gets its own
vector for fine-grained retrieval.

Chunking defaults
-----------------
CHUNK_SIZE     1200 chars  — comfortably under 512-token model limit
CHUNK_OVERLAP  150 chars   — sliding window so sentences don't get cut
"""
from __future__ import annotations

import base64
import re
from typing import TYPE_CHECKING

import httpx
import structlog

if TYPE_CHECKING:
    pass

logger = structlog.get_logger(__name__)

CHUNK_SIZE = 1200
CHUNK_OVERLAP = 150

# File size limit sent to Mistral OCR (50 MB)
MISTRAL_MAX_BYTES = 50 * 1024 * 1024

SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".md", ".docx", ".pptx"}
SUPPORTED_MIME = {
    "application/pdf",
    "text/plain",
    "text/markdown",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

async def extract_text(
    content: bytes,
    filename: str,
    content_type: str | None,
    mistral_api_key: str,
) -> str:
    """
    Route to the correct extractor and return the full extracted text.
    Raises ValueError for unsupported file types.
    Raises RuntimeError if Mistral OCR fails.
    """
    ext = _get_extension(filename, content_type)

    if ext in (".txt", ".md"):
        return _extract_txt(content)

    if ext == ".pdf":
        if not mistral_api_key:
            raise ValueError(
                "MISTRAL_API_KEY is not configured — cannot process PDF files."
            )
        if len(content) > MISTRAL_MAX_BYTES:
            raise ValueError(
                f"PDF is too large for Mistral OCR ({len(content) // 1024 // 1024} MB). "
                f"Maximum is {MISTRAL_MAX_BYTES // 1024 // 1024} MB."
            )
        return await _extract_pdf_mistral(content, filename, mistral_api_key)

    if ext == ".docx":
        return _extract_docx(content)

    if ext == ".pptx":
        return _extract_pptx(content)

    raise ValueError(
        f"Unsupported file type '{ext}'. "
        f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
    )


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def chunk_text(
    text: str,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> list[str]:
    """
    Split text into overlapping fixed-size chunks.

    Attempts to break at paragraph, then sentence, then word boundaries
    so chunks are coherent. Returns a single-element list if the text
    is shorter than chunk_size.
    """
    text = text.strip()
    if len(text) <= chunk_size:
        return [text] if text else []

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))

        # Try to break at a natural boundary in the last 25% of the chunk
        search_from = start + chunk_size * 3 // 4
        if end < len(text):
            for sep in ("\n\n", "\n", ". ", "! ", "? ", " "):
                pos = text.rfind(sep, search_from, end)
                if pos > start:
                    end = pos + len(sep)
                    break

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        if end >= len(text):
            break
        start = end - overlap

    return chunks


# ---------------------------------------------------------------------------
# Plain text
# ---------------------------------------------------------------------------

def _extract_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="replace").strip()


# ---------------------------------------------------------------------------
# PDF via Mistral OCR
# ---------------------------------------------------------------------------

async def _extract_pdf_mistral(
    pdf_bytes: bytes,
    filename: str,
    api_key: str,
) -> str:
    """
    Upload PDF to Mistral, run OCR, collect markdown per page, delete file.
    """
    headers_auth = {"Authorization": f"Bearer {api_key}"}
    file_id: str | None = None

    async with httpx.AsyncClient(timeout=180) as client:
        # 1. Upload the file
        upload_resp = await client.post(
            "https://api.mistral.ai/v1/files",
            headers=headers_auth,
            files={"file": (filename, pdf_bytes, "application/pdf")},
            data={"purpose": "ocr"},
        )
        upload_resp.raise_for_status()
        file_id = upload_resp.json()["id"]
        logger.debug("mistral_ocr.uploaded", file_id=file_id, filename=filename)

        # 2. Get signed URL (expiry=1 hour is the minimum)
        url_resp = await client.get(
            f"https://api.mistral.ai/v1/files/{file_id}/url",
            headers=headers_auth,
            params={"expiry": 1},
        )
        url_resp.raise_for_status()
        signed_url: str = url_resp.json()["url"]

        # 3. Run OCR
        ocr_resp = await client.post(
            "https://api.mistral.ai/v1/ocr",
            headers={**headers_auth, "Content-Type": "application/json"},
            json={
                "model": "mistral-ocr-latest",
                "document": {
                    "type": "document_url",
                    "document_url": signed_url,
                },
            },
        )
        ocr_resp.raise_for_status()
        pages: list[dict] = ocr_resp.json().get("pages", [])
        logger.debug("mistral_ocr.done", pages=len(pages))

        # 4. Cleanup — delete the uploaded file (best-effort)
        try:
            await client.delete(
                f"https://api.mistral.ai/v1/files/{file_id}",
                headers=headers_auth,
            )
        except Exception:
            pass

    # Combine page markdown, separated by clear page breaks
    page_texts: list[str] = []
    for page in pages:
        md = (page.get("markdown") or "").strip()
        if md:
            idx = page.get("index", len(page_texts))
            page_texts.append(f"<!-- page {idx + 1} -->\n{md}")

    return "\n\n".join(page_texts)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _extract_docx(content: bytes) -> str:
    try:
        from docx import Document  # type: ignore
    except ImportError:
        raise RuntimeError("python-docx is not installed. Run: pip install python-docx")

    from io import BytesIO
    doc = Document(BytesIO(content))
    parts: list[str] = []

    # Paragraphs (preserves headings)
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Mark headings so context is clear in the vector store
            if para.style.name.startswith("Heading"):
                parts.append(f"\n## {text}")
            else:
                parts.append(text)

    # Tables
    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                rows.append(row_text)
        if rows:
            parts.append("\n".join(rows))

    return "\n".join(parts).strip()


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    from io import BytesIO
    prs = Presentation(BytesIO(content))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"## Slide {i}"]

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if not text:
                continue
            # Title shapes first
            if shape.shape_type == 13 or "title" in (shape.name or "").lower():
                slide_parts.insert(1, f"**{text}**")
            else:
                slide_parts.append(text)

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"*Notes: {notes_text}*")

        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))

    return "\n\n".join(parts).strip()


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _get_extension(filename: str, content_type: str | None) -> str:
    """Derive a normalised file extension from filename or MIME type."""
    if filename:
        dot = filename.rfind(".")
        if dot >= 0:
            return filename[dot:].lower()

    # Fallback: map MIME type to extension
    mime_map = {
        "application/pdf": ".pdf",
        "text/plain": ".txt",
        "text/markdown": ".md",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    }
    return mime_map.get(content_type or "", "")
